#!/usr/bin/env python3
"""Inject speaker notes into a .pptx file's notes pane.

Minimal, standalone — no dependencies beyond python-pptx.
"""

import argparse
import json
import sys
from pathlib import Path


def ensure_pptx():
    try:
        import pptx  # noqa: F401
    except ImportError:
        sys.stderr.write(
            "python-pptx is required. Install with: pip install python-pptx\n"
        )
        raise SystemExit(2)


def set_notes_text(slide, text: str) -> None:
    """Set the notes text for a slide, creating a notes slide if needed."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    lines = text.split("\n")
    tf.text = lines[0] if lines else ""
    for extra in lines[1:]:
        p = tf.add_paragraph()
        p.text = extra


def inject(input_path: Path, output_path: Path, notes_map: dict[int, str]) -> list[str]:
    """Inject notes into PPTX. Returns a log list."""
    from pptx import Presentation

    prs = Presentation(str(input_path))
    total = len(prs.slides)
    log: list[str] = []

    for idx, slide in enumerate(prs.slides, start=1):
        if idx not in notes_map:
            log.append(f"slide {idx}: skipped (no entry)")
            continue
        new_text = notes_map[idx].strip()
        set_notes_text(slide, new_text)
        log.append(f"slide {idx}: injected ({len(new_text)} chars)")

    for extra in sorted(k for k in notes_map if k < 1 or k > total):
        log.append(f"WARNING: notes entry for slide {extra} ignored (deck has {total} slides)")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    log.append(f"saved: {output_path}")
    return log


def main(argv: list[str]) -> int:
    parser = argparse.ArgumentParser(
        description="Inject speaker notes into a .pptx file."
    )
    parser.add_argument("--input", required=True, type=Path, help="Input .pptx file")
    parser.add_argument("--output", required=True, type=Path, help="Output .pptx file")
    parser.add_argument("--notes", required=True, type=Path, help="JSON file: [{slide: N, notes: '...'}, ...]")
    args = parser.parse_args(argv)

    if not args.input.exists():
        sys.stderr.write(f"Input not found: {args.input}\n")
        return 1

    ensure_pptx()
    raw = json.loads(args.notes.read_text(encoding="utf-8"))
    notes_map = {int(e["slide"]): str(e["notes"]) for e in raw}

    for line in inject(args.input, args.output, notes_map):
        print(line)
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
